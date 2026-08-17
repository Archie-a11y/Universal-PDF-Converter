import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from tkinter.scrolledtext import ScrolledText
import os
import threading
import pythoncom
import win32com.client
from pdf2docx import Converter
from pdf2image import convert_from_path
from PIL import Image
import img2pdf
import fitz  # PyMuPDF
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF
from docx2pdf import convert

class PDFUniversalConverter:
    def __init__(self, root):
        self.root = root
        self.language = "en"  # 默认中文
        self.theme_mode = "light"  # 默认浅色主题

        # 支持的所有转换类型（中英文）
        self.conversion_types = {
            "pdf_to": {
                "word": ("PDF转Word", "PDF to Word"),
                "excel": ("PDF转Excel", "PDF to Excel"),
                "ppt": ("PDF转PPT", "PDF to PPT"),
                "pptx": ("PDF转PPTX", "PDF to PPTX"),
                "image": ("PDF转图片", "PDF to Images"),
                "text": ("PDF转文本", "PDF to Text")
            },
            "to_pdf": {
                "word": ("Word转PDF", "Word to PDF"),
                "excel": ("Excel转PDF", "Excel to PDF"),
                "ppt": ("PPT转PDF", "PPT to PDF"),
                "pptx": ("PPTX转PDF", "PPTX to PDF"),
                "image": ("图片转PDF", "Images to PDF"),
                "text": ("文本转PDF", "Text to PDF")
            }
        }

        # 完整翻译字典（包含所有界面元素）
        self.translations = {
            "zh": {
                # 窗口标题
                "title": "📄 全能PDF转换器",
                
                # 菜单栏
                "menu_language": "语言",
                "menu_theme": "主题",
                "menu_zh": "中文",
                "menu_en": "English",
                "menu_light": "浅色主题",
                "menu_dark": "深色主题",
                
                # 转换类型选择
                "select_conversion_type": "选择转换类型",
                "pdf_to": "PDF转换为",
                "to_pdf": "转换为PDF",
                
                # 文件选择
                "select_input": "选择输入文件",
                "select_output": "选择输出目录",
                "output_dir": "输出目录",
                "browse": "浏览",
                
                # 按钮
                "start_convert": "开始转换",
                "cancel": "取消",
                
                # 状态信息
                "status_ready": "准备就绪",
                "status_processing": "处理中...",
                "status_complete": "转换完成",
                "status_error": "转换失败",
                "converting": "正在转换: {}",
                
                # 日志
                "log_title": "操作日志",
                "file_selected": "已选择文件: {}",
                "output_selected": "输出目录: {}",
                
                # 消息
                "conversion_success": "转换成功! 文件保存在: {}",
                "conversion_failed": "转换失败: {}",
                "no_file_selected": "请先选择输入文件",
                "no_output_dir": "请先选择输出目录",
                "file_not_found": "文件不存在: {}",
                "invalid_file": "无效的文件格式",
                
                # 错误提示
                "error_title": "错误",
                "warning_title": "警告",
                "office_required": "需要安装Microsoft Office",
                "install_office": "请安装Microsoft Office以获得完整功能",
                "missing_dependency": "缺少依赖库: {}"
            },
            "en": {
                # Window title
                "title": "📄 Universal PDF Converter",
                
                # Menu items
                "menu_language": "Language",
                "menu_theme": "Theme",
                "menu_zh": "Chinese",
                "menu_en": "English",
                "menu_light": "Light Theme",
                "menu_dark": "Dark Theme",
                
                # Conversion type selection
                "select_conversion_type": "Select Conversion Type",
                "pdf_to": "PDF to",
                "to_pdf": "To PDF",
                
                # File selection
                "select_input": "Select Input File",
                "select_output": "Select Output Directory",
                "output_dir": "Output Directory",
                "browse": "Browse",
                
                # Buttons
                "start_convert": "Start Conversion",
                "cancel": "Cancel",
                
                # Status messages
                "status_ready": "Ready",
                "status_processing": "Processing...",
                "status_complete": "Conversion Complete",
                "status_error": "Conversion Failed",
                "converting": "Converting: {}",
                
                # Log
                "log_title": "Operation Log",
                "file_selected": "File selected: {}",
                "output_selected": "Output directory: {}",
                
                # Messages
                "conversion_success": "Conversion successful! File saved at: {}",
                "conversion_failed": "Conversion failed: {}",
                "no_file_selected": "Please select input file first",
                "no_output_dir": "Please select output directory first",
                "file_not_found": "File not found: {}",
                "invalid_file": "Invalid file format",
                
                # Error messages
                "error_title": "Error",
                "warning_title": "Warning",
                "office_required": "Microsoft Office required",
                "install_office": "Please install Microsoft Office for full functionality",
                "missing_dependency": "Missing dependency: {}"
            }
        }

        # 初始化变量
        self.input_path = ""
        self.output_dir = ""
        self.current_conversion_type = "pdf_to_word"
        
        # 初始化界面
        self.init_main_window()
        self.create_menu()
        self.setup_styles()
        self.build_gui()
        self.apply_theme()

    def _(self, text):
        """翻译文本"""
        return self.translations[self.language].get(text, text)

    def init_main_window(self):
        """初始化主窗口设置"""
        self.root.title(self._("title"))
        self.root.geometry("800x700")
        self.root.resizable(True, True)
        self.root.configure(bg=self.get_theme_colors()['root_bg'])

    def get_theme_colors(self):
        """获取当前主题颜色配置"""
        if self.theme_mode == "light":
            return {
                'root_bg': '#f5f5f5',
                'bg': '#ffffff',
                'fg': '#333333',
                'entry_bg': '#ffffff',
                'button_bg': '#4a90e2',
                'button_fg': '#ffffff',
                'text_fg': '#2a6eaf',
                'status_fg': '#666666',
                'menu_bg': '#f5f5f5',
                'menu_fg': '#333333',
                'menu_active_bg': '#e0e0e0',
                'menu_active_fg': '#000000',
                'frame_bg': '#f5f5f5',
                'text_bg': '#ffffff',
                'text_fg': '#000000',
                'combobox_bg': '#ffffff',
                'separator_color': '#cccccc'
            }
        else:
            return {
                'root_bg': '#2d2d2d',
                'bg': '#3d3d3d',
                'fg': '#ffffff',
                'entry_bg': '#3d3d3d',
                'button_bg': '#1a73e8',
                'button_fg': '#ffffff',
                'text_fg': '#5a9ed9',
                'status_fg': '#aaaaaa',
                'menu_bg': '#2d2d2d',
                'menu_fg': '#ffffff',
                'menu_active_bg': '#1a1a1a',
                'menu_active_fg': '#ffffff',
                'frame_bg': '#2d2d2d',
                'text_bg': '#3d3d3d',
                'text_fg': '#ffffff',
                'combobox_bg': '#3d3d3d',
                'separator_color': '#555555'
            }

    def create_menu(self):
        """创建菜单栏（完整中英文支持）"""
        colors = self.get_theme_colors()
        self.menubar = tk.Menu(self.root, 
                              bg=colors['menu_bg'], 
                              fg=colors['menu_fg'],
                              activebackground=colors['menu_active_bg'],
                              activeforeground=colors['menu_active_fg'])
        
        # 语言菜单
        self.language_menu = tk.Menu(self.menubar, 
                                   tearoff=0,
                                   bg=colors['menu_bg'],
                                   fg=colors['menu_fg'],
                                   activebackground=colors['menu_active_bg'],
                                   activeforeground=colors['menu_active_fg'])
        self.language_menu.add_command(label=self._("menu_zh"), 
                                     command=lambda: self.change_language("zh"))
        self.language_menu.add_command(label=self._("menu_en"), 
                                     command=lambda: self.change_language("en"))
        self.menubar.add_cascade(label=self._("menu_language"), 
                                menu=self.language_menu)
        
        # 主题菜单
        self.theme_menu = tk.Menu(self.menubar, 
                                tearoff=0,
                                bg=colors['menu_bg'],
                                fg=colors['menu_fg'],
                                activebackground=colors['menu_active_bg'],
                                activeforeground=colors['menu_active_fg'])
        self.theme_menu.add_command(label=self._("menu_light"), 
                                  command=lambda: self.change_theme("light"))
        self.theme_menu.add_command(label=self._("menu_dark"), 
                                  command=lambda: self.change_theme("dark"))
        self.menubar.add_cascade(label=self._("menu_theme"), 
                               menu=self.theme_menu)
        
        self.root.config(menu=self.menubar)

    def setup_styles(self):
        """设置UI样式"""
        colors = self.get_theme_colors()
        self.style = ttk.Style()
        self.style.theme_use('clam')
        
        # 配置全局样式
        self.style.configure('.', 
                           background=colors['frame_bg'],
                           foreground=colors['fg'],
                           font=('Arial', 11))
        
        # 配置各种组件样式
        self.style.configure('TFrame', background=colors['frame_bg'])
        self.style.configure('TLabel', 
                           background=colors['frame_bg'],
                           foreground=colors['fg'])
        self.style.configure('TEntry', 
                           fieldbackground=colors['entry_bg'],
                           foreground=colors['fg'],
                           insertcolor=colors['fg'],
                           padding=8)
        self.style.configure('TCombobox', 
                           fieldbackground=colors['combobox_bg'],
                           foreground=colors['fg'],
                           padding=8)
        self.style.configure('TButton', 
                           background=colors['button_bg'],
                           foreground=colors['button_fg'],
                           font=('Arial', 12, 'bold'),
                           padding=10)
        self.style.map('TButton',
                      background=[('active', '#3a7ebf'), ('pressed', '#2a6eaf')])
        self.style.configure('TLabelframe', 
                           background=colors['frame_bg'],
                           foreground=colors['fg'])
        self.style.configure('Horizontal.TSeparator',
                           background=colors['separator_color'])

    def build_gui(self):
        """构建用户界面（完整中英文支持）"""
        colors = self.get_theme_colors()
        
        # 主容器
        self.main_frame = ttk.Frame(self.root)
        self.main_frame.pack(expand=True, fill=tk.BOTH, padx=20, pady=20)
        
        # 1. 转换类型选择区域
        self.type_frame = ttk.LabelFrame(self.main_frame, text=self._("select_conversion_type"))
        self.type_frame.pack(fill=tk.X, pady=10, padx=5)
        
        # 转换方向选择 (PDF转其他/其他转PDF)
        self.direction_frame = ttk.Frame(self.type_frame)
        self.direction_frame.pack(fill=tk.X, pady=5)
        
        self.conversion_direction = tk.StringVar(value="pdf_to")
        ttk.Radiobutton(self.direction_frame, 
                       text=self._("pdf_to"),
                       variable=self.conversion_direction,
                       value="pdf_to",
                       command=self.update_conversion_types).pack(side=tk.LEFT, padx=10)
        ttk.Radiobutton(self.direction_frame, 
                       text=self._("to_pdf"),
                       variable=self.conversion_direction,
                       value="to_pdf",
                       command=self.update_conversion_types).pack(side=tk.LEFT, padx=10)
        
        # 具体转换类型选择
        self.conversion_type = tk.StringVar()
        self.type_combobox = ttk.Combobox(self.type_frame, 
                                         textvariable=self.conversion_type,
                                         state="readonly")
        self.type_combobox.pack(fill=tk.X, pady=5, padx=5)
        self.update_conversion_types()
        
        # 2. 文件选择区域
        self.file_frame = ttk.LabelFrame(self.main_frame, text=self._("select_input"))
        self.file_frame.pack(fill=tk.X, pady=10, padx=5)
        
        # 输入文件选择
        self.input_frame = ttk.Frame(self.file_frame)
        self.input_frame.pack(fill=tk.X, pady=5)
        
        self.input_path_var = tk.StringVar()
        ttk.Entry(self.input_frame, 
                 textvariable=self.input_path_var, 
                 state='readonly').pack(side=tk.LEFT, expand=True, fill=tk.X, padx=5)
        ttk.Button(self.input_frame, 
                  text=self._("select_input"),
                  command=self.select_input_file).pack(side=tk.LEFT, padx=5)
        
        # 输出目录选择
        self.output_frame = ttk.Frame(self.file_frame)
        self.output_frame.pack(fill=tk.X, pady=5)
        
        self.output_dir_var = tk.StringVar()
        ttk.Entry(self.output_frame, 
                 textvariable=self.output_dir_var, 
                 state='readonly').pack(side=tk.LEFT, expand=True, fill=tk.X, padx=5)
        ttk.Button(self.output_frame, 
                  text=self._("select_output"),
                  command=self.select_output_dir).pack(side=tk.LEFT, padx=5)
        
        # 3. 转换按钮区域
        self.button_frame = ttk.Frame(self.main_frame)
        self.button_frame.pack(fill=tk.X, pady=10)
        
        ttk.Button(self.button_frame, 
                 text=self._("start_convert"),
                 command=self.start_conversion).pack(side=tk.LEFT, padx=5)
        
        ttk.Button(self.button_frame,
                 text=self._("cancel"),
                 command=self.cancel_conversion).pack(side=tk.LEFT, padx=5)
        
        # 分隔线
        ttk.Separator(self.main_frame, orient='horizontal').pack(fill=tk.X, pady=10)
        
        # 4. 日志区域
        self.log_frame = ttk.LabelFrame(self.main_frame, text=self._("log_title"))
        self.log_frame.pack(expand=True, fill=tk.BOTH, pady=10)
        
        self.log_text = ScrolledText(self.log_frame, 
                                   wrap=tk.WORD, 
                                   state='disabled',
                                   font=('Consolas', 10))
        self.log_text.pack(expand=True, fill=tk.BOTH, padx=5, pady=5)
        
        # 5. 状态栏
        self.status_frame = ttk.Frame(self.main_frame)
        self.status_frame.pack(fill=tk.X, pady=(0, 10))
        
        self.status_var = tk.StringVar()
        self.status_var.set(self._("status_ready"))
        ttk.Label(self.status_frame, 
                textvariable=self.status_var,
                font=('Arial', 9)).pack(side=tk.LEFT)

    def update_conversion_types(self):
        """更新可选的转换类型（完整中英文支持）"""
        direction = self.conversion_direction.get()
        types = self.conversion_types[direction]
        
        # 根据语言获取显示文本
        display_texts = [types[t][0] if self.language == "zh" else types[t][1] 
                        for t in types]
        
        self.type_combobox['values'] = display_texts
        self.type_combobox.current(0)  # 默认选择第一个

    def select_input_file(self):
        """选择输入文件"""
        direction = self.conversion_direction.get()
        filetypes = []
        
        # 根据转换方向设置文件类型过滤器
        if direction == "pdf_to":
            filetypes = [("PDF Files", "*.pdf"), ("All Files", "*.*")]
        else:  # to_pdf
            current_type = self.get_current_conversion_key()
            if current_type == "word":
                filetypes = [("Word Files", "*.docx *.doc"), ("All Files", "*.*")]
            elif current_type == "excel":
                filetypes = [("Excel Files", "*.xlsx *.xls"), ("All Files", "*.*")]
            elif current_type == "ppt":
                filetypes = [("PowerPoint Files", "*.pptx *.ppt"), ("All Files", "*.*")]
            elif current_type == "pptx":
                filetypes = [("PowerPoint Files", "*.pptx"), ("All Files", "*.*")]
            elif current_type == "image":
                filetypes = [("Image Files", "*.jpg *.jpeg *.png *.bmp"), ("All Files", "*.*")]
            elif current_type == "text":
                filetypes = [("Text Files", "*.txt"), ("All Files", "*.*")]
        
        file_path = filedialog.askopenfilename(
            title=self._("select_input"),
            filetypes=filetypes
        )
        
        if file_path:
            self.input_path = file_path
            self.input_path_var.set(file_path)
            self.log_message(self._("file_selected").format(file_path))

    def select_output_dir(self):
        """选择输出目录"""
        dir_path = filedialog.askdirectory(title=self._("output_dir"))
        if dir_path:
            self.output_dir = dir_path
            self.output_dir_var.set(dir_path)
            self.log_message(self._("output_selected").format(dir_path))

    def get_current_conversion_key(self):
        """获取当前选择的转换类型key"""
        direction = self.conversion_direction.get()
        display_text = self.conversion_type.get()
        
        # 在字典中查找匹配的key
        for key, (zh_text, en_text) in self.conversion_types[direction].items():
            if display_text == (zh_text if self.language == "zh" else en_text):
                return key
        return "word"  # 默认返回word

    def start_conversion(self):
        """开始转换"""
        if not self.input_path:
            messagebox.showwarning(self._("warning_title"), self._("no_file_selected"))
            return
        
        if not os.path.exists(self.input_path):
            messagebox.showwarning(self._("warning_title"), 
                                 self._("file_not_found").format(self.input_path))
            return
        
        if not self.output_dir:
            messagebox.showwarning(self._("warning_title"), self._("no_output_dir"))
            return
        
        self.status_var.set(self._("status_processing"))
        self.log_message(self._("converting").format(os.path.basename(self.input_path)))
        
        # 在后台线程中执行转换
        threading.Thread(target=self.perform_conversion, daemon=True).start()

    def cancel_conversion(self):
        """取消转换"""
        self.log_message("Conversion cancelled by user")
        self.status_var.set(self._("status_ready"))

    def perform_conversion(self):
        """执行实际的转换操作"""
        try:
            direction = self.conversion_direction.get()
            conv_type = self.get_current_conversion_key()
            input_file = self.input_path
            output_dir = self.output_dir
            
            # 获取输出文件名
            base_name = os.path.splitext(os.path.basename(input_file))[0]
            
            if direction == "pdf_to":
                # PDF转其他格式
                if conv_type == "word":
                    output_file = os.path.join(output_dir, f"{base_name}.docx")
                    self.pdf_to_word(input_file, output_file)
                elif conv_type == "excel":
                    output_file = os.path.join(output_dir, f"{base_name}.xlsx")
                    self.pdf_to_excel(input_file, output_file)
                elif conv_type == "ppt":
                    output_file = os.path.join(output_dir, f"{base_name}.ppt")
                    self.pdf_to_ppt(input_file, output_file)
                elif conv_type == "pptx":
                    output_file = os.path.join(output_dir, f"{base_name}.pptx")
                    self.pdf_to_pptx(input_file, output_file)
                elif conv_type == "image":
                    output_file = os.path.join(output_dir, f"{base_name}")
                    self.pdf_to_image(input_file, output_file)
                elif conv_type == "text":
                    output_file = os.path.join(output_dir, f"{base_name}.txt")
                    self.pdf_to_text(input_file, output_file)
            else:
                # 其他格式转PDF
                if conv_type == "word":
                    output_file = os.path.join(output_dir, f"{base_name}.pdf")
                    self.word_to_pdf(input_file, output_file)
                elif conv_type == "excel":
                    output_file = os.path.join(output_dir, f"{base_name}.pdf")
                    self.excel_to_pdf(input_file, output_file)
                elif conv_type == "ppt":
                    output_file = os.path.join(output_dir, f"{base_name}.pdf")
                    self.ppt_to_pdf(input_file, output_file)
                elif conv_type == "pptx":
                    output_file = os.path.join(output_dir, f"{base_name}.pdf")
                    self.pptx_to_pdf(input_file, output_file)
                elif conv_type == "image":
                    output_file = os.path.join(output_dir, f"{base_name}.pdf")
                    self.image_to_pdf(input_file, output_file)
                elif conv_type == "text":
                    output_file = os.path.join(output_dir, f"{base_name}.pdf")
                    self.text_to_pdf(input_file, output_file)
            
            self.log_message(self._("conversion_success").format(output_file))
            self.status_var.set(self._("status_complete"))
            
        except Exception as e:
            self.log_message(self._("conversion_failed").format(str(e)))
            self.status_var.set(self._("status_error"))
            messagebox.showerror(self._("error_title"), str(e))

    # ==================== 转换方法实现 ====================

    def pdf_to_word(self, pdf_path, output_path):
        """PDF转Word"""
        try:
            cv = Converter(pdf_path)
            cv.convert(output_path, start=0, end=None)
            cv.close()
        except Exception as e:
            raise Exception(f"PDF to Word conversion failed: {str(e)}")

    def pdf_to_excel(self, pdf_path, output_path):
        """PDF转Excel"""
        try:
            doc = fitz.open(pdf_path)
            text = ""
            for page in doc:
                text += page.get_text()
            
            # 简单处理：按行分割，制表符分隔
            lines = text.split('\n')
            data = []
            for line in lines:
                if line.strip():
                    # 尝试用多种分隔符分割
                    parts = line.replace('\t', '|').replace('  ', '|').split('|')
                    parts = [p.strip() for p in parts if p.strip()]
                    if parts:
                        data.append(parts)
            
            # 确保所有行有相同数量的列
            max_cols = max(len(row) for row in data) if data else 0
            for row in data:
                while len(row) < max_cols:
                    row.append("")
            
            df = pd.DataFrame(data)
            df.to_excel(output_path, index=False, header=False)
        except Exception as e:
            raise Exception(f"PDF to Excel conversion failed: {str(e)}\n{self._('office_required')}")

    def pdf_to_ppt(self, pdf_path, output_path):
        """PDF转PPT"""
        try:
            images = convert_from_path(pdf_path)
            prs = Presentation()
            
            # 使用16:9的幻灯片布局
            prs.slide_width = Inches(16).emu
            prs.slide_height = Inches(9).emu
            
            for i, img in enumerate(images):
                # 创建新幻灯片
                blank_slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # 临时保存图片
                img_path = f"temp_pdf_page_{i}.jpg"
                img.save(img_path, "JPEG")
                
                # 添加图片到幻灯片并居中
                slide.shapes.add_picture(img_path, 0, 0, 
                                       width=prs.slide_width, 
                                       height=prs.slide_height)
                
                # 删除临时文件
                os.remove(img_path)
            
            prs.save(output_path)
        except Exception as e:
            raise Exception(f"PDF to PPT conversion failed: {str(e)}")

    def pdf_to_pptx(self, pdf_path, output_path):
        """PDF转PPTX"""
        try:
            doc = fitz.open(pdf_path)
            prs = Presentation()

            # 使用16:9的幻灯片布局
            prs.slide_width = Inches(16).emu
            prs.slide_height = Inches(9).emu

            for page in doc:
                # 将 PDF 页面转为临时图片
                img_path = f"temp_page_{page.number}.png"
                pix = page.get_pixmap(dpi=200)  # 设置分辨率
                pix.save(img_path)

                # 创建新幻灯片
                blank_slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(blank_slide_layout)

                # 添加图片到幻灯片并居中
                slide.shapes.add_picture(img_path, 0, 0, 
                                         width=prs.slide_width, 
                                         height=prs.slide_height)

                # 删除临时文件
                os.remove(img_path)

            prs.save(output_path)
        except Exception as e:
            raise Exception(f"PDF to PPTX conversion failed: {str(e)}")

    def pdf_to_image(self, pdf_path, output_prefix):
        """PDF转图片"""
        try:
            images = convert_from_path(pdf_path)
            for i, image in enumerate(images):
                image.save(f"{output_prefix}_page_{i+1}.jpg", "JPEG")
        except Exception as e:
            raise Exception(f"PDF to Images conversion failed: {str(e)}")

    def pdf_to_text(self, pdf_path, output_path):
        """PDF转文本"""
        try:
            doc = fitz.open(pdf_path)
            with open(output_path, "w", encoding="utf-8") as f:
                for page in doc:
                    text = page.get_text()
                    # 清理文本格式
                    text = text.replace('\r\n', '\n').replace('\r', '\n')
                    lines = [line.strip() for line in text.split('\n') if line.strip()]
                    f.write('\n'.join(lines) + '\n\n')
        except Exception as e:
            raise Exception(f"PDF to Text conversion failed: {str(e)}")

    def word_to_pdf(self, word_path, output_path):
        """Word转PDF"""
        try:
            # 使用docx2pdf
            convert(word_path, output_path)
        except Exception as e:
            # 如果docx2pdf失败，尝试使用win32com
            try:
                pythoncom.CoInitialize()
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(os.path.abspath(word_path))
                doc.SaveAs(os.path.abspath(output_path), FileFormat=17)  # 17是PDF格式
                doc.Close()
                word.Quit()
            except Exception as win32_error:
                raise Exception(f"{self._('office_required')}: {str(win32_error)}")

    def excel_to_pdf(self, excel_path, output_path):
        """Excel转PDF"""
        try:
            pythoncom.CoInitialize()
            excel = win32com.client.Dispatch("Excel.Application")
            excel.Visible = False
            workbook = excel.Workbooks.Open(os.path.abspath(excel_path))
            
            # 设置打印区域为整个工作表
            worksheet = workbook.ActiveSheet
            worksheet.PageSetup.Zoom = False
            worksheet.PageSetup.FitToPagesWide = 1
            worksheet.PageSetup.FitToPagesTall = False
            
            # 导出为PDF
            workbook.ExportAsFixedFormat(0, os.path.abspath(output_path))  # 0是PDF格式
            workbook.Close()
            excel.Quit()
        except Exception as e:
            raise Exception(f"{self._('office_required')}: {str(e)}")

    def ppt_to_pdf(self, ppt_path, output_path):
        """PPT转PDF"""
        try:
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = 1
            presentation = powerpoint.Presentations.Open(os.path.abspath(ppt_path))
            presentation.ExportAsFixedFormat(
                os.path.abspath(output_path), 
                2,  # ppFixedFormatTypePDF
                PrintRange=None
            )
            presentation.Close()
            powerpoint.Quit()
        except Exception as e:
            raise Exception(f"{self._('office_required')}: {str(e)}")

    def pptx_to_pdf(self, pptx_path, output_path):
        """PPTX转PDF"""
        try:
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = 1
            presentation = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
            presentation.ExportAsFixedFormat(
                os.path.abspath(output_path), 
                2,  # ppFixedFormatTypePDF
                PrintRange=None
            )
            presentation.Close()
            powerpoint.Quit()
        except Exception as e:
            raise Exception(f"{self._('office_required')}: {str(e)}")

    def image_to_pdf(self, image_path, output_path):
        """图片转PDF"""
        try:
            # 验证图片文件
            try:
                with Image.open(image_path) as img:
                    img.verify()
            except Exception as e:
                raise Exception(f"Invalid image file: {str(e)}")
            
            with open(output_path, "wb") as f:
                f.write(img2pdf.convert(image_path))
        except Exception as e:
            raise Exception(f"Images to PDF conversion failed: {str(e)}")

    def text_to_pdf(self, text_path, output_path):
        """文本转PDF"""
        try:
            # 读取文本内容
            with open(text_path, "r", encoding="utf-8") as f:
                text_content = f.read()
            
            # 创建PDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            
            # 添加文本
            pdf.multi_cell(0, 10, txt=text_content)
            
            # 保存PDF
            pdf.output(output_path)
        except Exception as e:
            raise Exception(f"Text to PDF conversion failed: {str(e)}")

    def log_message(self, message):
        """在日志区域添加消息"""
        self.log_text.config(state='normal')
        self.log_text.insert(tk.END, message + "\n")
        self.log_text.see(tk.END)
        self.log_text.config(state='disabled')

    def change_language(self, lang):
        """切换语言并更新整个UI"""
        self.language = lang
        self.update_ui_text()
        self.root.title(self._("title"))
        self.update_conversion_types()
        self.create_menu()

    def update_ui_text(self):
        """更新所有UI文本（完整中英文支持）"""
        # 更新框架标题
        self.type_frame.config(text=self._("select_conversion_type"))
        self.file_frame.config(text=self._("select_input"))
        self.log_frame.config(text=self._("log_title"))
        
        # 更新按钮文本
        for child in self.button_frame.winfo_children():
            if isinstance(child, ttk.Button):
                if "start_convert" in str(child.cget('command')):
                    child.config(text=self._("start_convert"))
                elif "cancel" in str(child.cget('command')):
                    child.config(text=self._("cancel"))
        
        # 更新浏览按钮文本
        for frame in [self.input_frame, self.output_frame]:
            for child in frame.winfo_children():
                if isinstance(child, ttk.Button):
                    child.config(text=self._("browse"))
        
        # 更新单选按钮文本
        for child in self.direction_frame.winfo_children():
            if isinstance(child, ttk.Radiobutton):
                if child.cget('value') == "pdf_to":
                    child.config(text=self._("pdf_to"))
                else:
                    child.config(text=self._("to_pdf"))
        
        # 更新状态栏
        current_status = self.status_var.get()
        if current_status in [self.translations['zh']['status_ready'],
                             self.translations['en']['status_ready'],
                             self.translations['zh']['status_processing'],
                             self.translations['en']['status_processing'],
                             self.translations['zh']['status_complete'],
                             self.translations['en']['status_complete'],
                             self.translations['zh']['status_error'],
                             self.translations['en']['status_error']]:
            self.status_var.set(self._("status_ready"))

    def change_theme(self, theme):
        """切换主题并更新整个UI"""
        self.theme_mode = theme
        self.apply_theme()
        self.root.configure(bg=self.get_theme_colors()['root_bg'])
        self.create_menu()

    def apply_theme(self):
        """应用当前主题到所有组件"""
        colors = self.get_theme_colors()
        
        # 重新配置样式
        self.setup_styles()
        
        # 更新日志区域颜色
        self.log_text.config(
            bg=colors['text_bg'],
            fg=colors['text_fg'],
            insertbackground=colors['text_fg']
        )


if __name__ == "__main__":
    # 初始化COM库
    pythoncom.CoInitialize()

    root = tk.Tk()
    try:
        app = PDFUniversalConverter(root)
        root.mainloop()
    finally:
        # 确保程序退出时清理COM库
        pythoncom.CoUninitialize()